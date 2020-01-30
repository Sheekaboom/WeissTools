'''
@brief this module is to let us quickly build some nice graphs. Graphviz works 
really well for certain scenarios but does not give as much flexibility as one
would hope with layout. This library will hopefully eventually implement all of 
graphviz functionality but extended. For now it will allow creation of flow-charts
in a gridded fasion. This will try to follow some keywords for graphviz

@
@note svg not yet implemented https://svgwrite.readthedocs.io/en/

@author ajw
@date 1-2019
'''

import numpy as np
import svgwrite
from pptx import Presentation

COLOR_DICT = {
    'black'  :  '#000000',
    'white'  :  '#FFFFFF',
    'grey'   :  '#7F7F7F',
    'red'    :  '#FF0000',
    'green'  :  '#00FF00',
    'blue'   :  '#0000FF',
    }

class DEFAULT_ID:
    '''@brief class for getting a default id'''
    _VAL = 1000000 
    @classmethod
    def next(cls):
        '''@brief move to next id'''
        cls._VAL+=1
    @classmethod
    def __str__(cls):
        '''@brief get the id as a string'''
        return '##{}##'.format(cls._VAL)
    @classmethod
    def get(cls,inc=True):
        '''@brief get the id and increment'''
        rv = cls.__str__()
        if inc:
            cls.next()
        return rv

def rgb2hex(*args):
    '''
    @brief convert rgb colors to hex
    @param[in] args - can be a list of rgb values (e.g. [r,g,b]) or separate args
        (e.g. r,g,b)
    @cite https://stackoverflow.com/questions/3380726/converting-a-rgb-color-tuple-to-a-six-digit-code-in-python
    '''
    if len(args)>1: #then assume its r,g,b
        rgb_vals = np.array(args)
    else:
        rgb_vals = np.array(args[0]) #otherwise get [r,g,b]
    
    hex_str = '#{:02X}{:02X}{:02X}'.format(*tuple(rgb_vals))
    return hex_str

class GraphObj(dict):
    '''
    @brief this is a superclass for graph objects
    @param[in] id - id or name of object for referencing. if not given, autogenerate
    @param[in/OPT] kwargs - kwargs for dict.__init__
    @note this will mostly be for setting properties
    '''
    def __init__(self,id=None,**kwargs):
        if id is None:
            id = DEFAULT_ID.get()
        self.id = id
        
        super().__init__(**kwargs)
        #first lets set our default params
        default_dict = {
            'style'    : None,
            'color'    : COLOR_DICT['black'],
            'bgcolor'  : COLOR_DICT['white'],
            'text'     : '',
            'location' : [0,0], #location to place (probably xy on a grid)
            'drawing'  : dict(), #info for drawing later
            }
        for k,v in default_dict.items():
            if k not in self: #set the defaults if they werent already set
                self[k] = v
        
class Node(GraphObj):
    '''
    @brief class for a node in the graph
    @param[in] id - id or name of object for referencing
    @param[in/OPT] kwargs - kwargs for dict.__init__
    '''
    def __init__(self,id,**kwargs):
        #self['edges'] = set() #input edges to the node
        super().__init__(id=id,**kwargs)
        
class Edge(GraphObj):
    '''
    @brief class for a node in the graph
    @param[in] args - nodes to connect
    @param[in] id - id or name of object for referencing
    @param[in/OPT] kwargs - kwargs for dict.__init__
    '''
    def __init__(self,*args,**kwargs):
        self['nodes'] = set(args) #nodes the edge is connected to
        super().__init__(**kwargs)
        
class Graph(GraphObj):
    '''
    @brief instantiate a graph object (or subgraph)
    @param[in] id - id or name of graph for referencing
    @param[in/OPT] kwargs - keyword arguments for super.__init__
    '''
    def __init__(self,id=None,**kwargs):
        if id is None:
            id = DEFAULT_ID.get()
        #init with no edges or nodes
        self['rankdir'] = 'LR'
        self['edges'] = {}
        self['nodes'] = {}
        self['subgraphs'] = {}
        #init superclass
        super().__init__(id,**kwargs)
        
    def add_node(self,id,**kwargs):
        '''
        @brief add a node to the graph
        @param[in] id - id to reference the node
        @param[in] args - arguments passed to Node() 
        @param[in] kwargs - keyword args passed to Node()
        @return Object of type Node
        '''
        node = Node(id,**kwargs)
        self['nodes'][id] = node
        return node
    
    def add_edge(self,node_a,node_b,**kwargs):
        '''
        @brief add an edge between two nodes or subgraphs
        @note the nodes must exist in the current graph or subgraphs
        @param[in] node_a - node to start from
        @param[in] node_b - node to end at 
        @param[in/OPT] kwargs - passed to Edge init
        @return Object of type Edge
        '''
        edge = Edge(node_a,node_b,**kwargs)
        id = edge.id
        self['edges'][id] = edge 
        return edge
    
    @property
    def nodes(self): return self['nodes']
    @property
    def edges(self): return self['edges']
    

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
    
class GraphDrawer:
    '''
    @brief class to draw a PyGraph.Graph
    @note this only currently works with pptx
    '''
    def __init__():
        
        self.prs = Presnetation
        self.blank_layout = self.prs.slide_layouts[6] #blank slide
        self.slide = self.prs.slides.add_slide(self.blank_layout)
    
    def draw_node(node): #create a group and draw an svg node
       '''
       @brief draw a node in svg format with text and size
       @param[in] node - node to draw
       @note this updates the 'drawing' item in the node
       '''
       drawing_node_dict = {
           'x_size'  : 1,
           'y_size'  : 1,
           'x_center': 0.5,
           'y_center': 0.5,
           }
       node = self.dwg.g(id=node.id)
       text = self.dwg.text
       rect = self.dwg.rect()
       rect = node.add(self.dwg.rect())
       
       
       
    #def draw_grid(): 
    #    @brief draw a grid with given size to put plots in
        
        
        
import unittest
class TestPyGraph(unittest.TestCase):
    '''@brief test some PyGraph funcionality'''
    
    def test_rgb2hex(self):
        '''@brief test rgb2hex program'''
        self.assertEqual(rgb2hex(0,0,0),'#000000')
        self.assertEqual(rgb2hex([255,255,255]),'#FFFFFF')
        
    
if __name__=='__main__':
    
    unittest.main() 
    
    G = Graph()
    G.add_node('A')
    G.add_node('B',location=[0,1])
    G.add_edge('A','B')
    
    
    
    

    